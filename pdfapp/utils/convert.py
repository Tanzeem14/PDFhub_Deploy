


import os

from django.conf import settings


def convert_pdf_to_word(pdf_file, output_file):
    from pdf2docx import Converter
    cv = Converter(pdf_file)
    cv.convert(output_file, start=0, end=None)
    cv.close()



def convert_pdf_to_images(pdf_file, output_dir):
    from pdf2image import convert_from_path

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    images = convert_from_path(pdf_file, poppler_path=getattr(settings, "POPPLER_PATH", r"C:\Users\lenovo\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin"))
    image_paths = []

    for i, image in enumerate(images):
        image_path = os.path.join(output_dir, f"page_{i + 1}.jpg")
        image.save(image_path, "JPEG")
        image_paths.append(image_path)

    return image_paths


def convert_pdf_to_pptx(pdf_file, output_file):
    
    from pptx import Presentation
    from pdf2image import convert_from_path
    images = convert_from_path(pdf_file, poppler_path=getattr(settings, "POPPLER_PATH", r"D:\PDFAPPLIBRARY\POPPLER\poppler-24.08.0\Library\bin"))
    presentation = Presentation()

    for image in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        temp_img = "temp_image.jpg"
        image.save(temp_img, "JPEG")
        slide.shapes.add_picture(temp_img, 0, 0, width=presentation.slide_width)
        os.remove(temp_img)

    presentation.save(output_file)
